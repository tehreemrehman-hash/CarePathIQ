from pptx import Presentation

def extract_text_from_pptx(file_path):
    """
    Extracts text from a .pptx file.
    """
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\\n".join(text_runs)
    except Exception as e:
        return f"An error occurred: {e}"

if __name__ == "__main__":
    pptx_file = "Course/SBM_presentation_6_2.pptx"
    extracted_text = extract_text_from_pptx(pptx_file)
    print(extracted_text)
